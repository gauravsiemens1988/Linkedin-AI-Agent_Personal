from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation(
    slides_data,
    image_folder="drafts/images",
    logo_path="assets/logo_watermark.png"
):
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title + Content (MOST STABLE)

    # -----------------------------
    # RESOLVE PATHS (CRITICAL)
    # -----------------------------
    base_dir = os.getcwd()
    image_folder_path = os.path.join(base_dir, image_folder)

    print("DEBUG: base_dir =", base_dir)
    print("DEBUG: image_folder_path =", image_folder_path)

    if os.path.exists(image_folder_path):
        print("DEBUG: image folder exists")
        print("DEBUG: files in image folder =", os.listdir(image_folder_path))
    else:
        print("DEBUG: image folder DOES NOT EXIST")

    # -----------------------------
    # PICK FIRST IMAGE
    # -----------------------------
    image_path = None
    if os.path.exists(image_folder_path):
        for f in sorted(os.listdir(image_folder_path)):
            if f.lower().endswith(".png"):
                image_path = os.path.join(image_folder_path, f)
                break

    print("DEBUG: resolved image_path =", image_path)

    # -----------------------------
    # CREATE SLIDES
    # -----------------------------
    for idx, slide_data in enumerate(slides_data, start=1):
        slide = prs.slides.add_slide(layout)

        # -----------------------------
        # TITLE
        # -----------------------------
        title = slide.shapes.title
        title.text = slide_data["title"]
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.bold = True

        # -----------------------------
        # CONTENT PLACEHOLDER
        # -----------------------------
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()

        # -----------------------------
        # IMAGE INSIDE PLACEHOLDER (SAFE)
        # -----------------------------
        if image_path and os.path.exists(image_path):
            content.insert_picture(image_path)
            print(f"DEBUG: Image inserted on slide {idx}")
        else:
            tf.text = "⚠️ Image missing (not found in runner)"
            print(f"DEBUG: Image NOT inserted on slide {idx}")

        # -----------------------------
        # BULLET POINTS (LIMITED)
        # -----------------------------
        for point in slide_data["points"][:4]:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.level = 1

        # -----------------------------
        # LOGO WATERMARK (OPTIONAL)
        # -----------------------------
        logo_full = os.path.join(base_dir, logo_path)
        if os.path.exists(logo_full):
            slide.shapes.add_picture(
                logo_full,
                prs.slide_width - Inches(1.8),
                prs.slide_height - Inches(0.9),
                width=Inches(1.4)
            )

    # -----------------------------
    # SAVE PRESENTATION
    # -----------------------------
    os.makedirs("drafts", exist_ok=True)
    output_path = os.path.join("drafts", "linkedin_carousel.pptx")
    prs.save(output_path)

    print("DEBUG: Presentation saved at", output_path)

    return output_path



